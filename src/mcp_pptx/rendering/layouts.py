"""Layout management for PowerPoint presentations."""

import logging
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.slide import SlideLayout

from ..models.deck_spec import LayoutType

logger = logging.getLogger(__name__)


class LayoutManager:
    """Manages slide layouts and fallbacks."""

    # Mapping of our layout types to typical PowerPoint layout indices
    LAYOUT_MAPPINGS = {
        LayoutType.TITLE: 0,           # Title slide
        LayoutType.TITLE_CONTENT: 1,   # Title and content
        LayoutType.SECTION: 2,         # Section header
        LayoutType.TWO_COL: 3,         # Two content columns
        LayoutType.IMAGE_FOCUS: 1,     # Fallback to title and content
        LayoutType.TABLE: 1,           # Fallback to title and content
        LayoutType.CHART: 1,           # Fallback to title and content
        LayoutType.BLANK: 6,           # Blank slide
    }

    def get_layout(self, prs: Presentation, layout_type: LayoutType) -> Optional[SlideLayout]:
        """Get slide layout by type with fallback."""
        try:
            # Get preferred layout index
            layout_index = self.LAYOUT_MAPPINGS.get(layout_type)
            
            if layout_index is not None and layout_index < len(prs.slide_layouts):
                return prs.slide_layouts[layout_index]
            
            # Try to find by name (common PowerPoint layout names)
            layout_names = {
                LayoutType.TITLE: ["Title Slide", "Title Only"],
                LayoutType.TITLE_CONTENT: ["Title and Content", "Content with Caption"],
                LayoutType.SECTION: ["Section Header", "Title Only"],
                LayoutType.TWO_COL: ["Two Content", "Comparison"],
                LayoutType.BLANK: ["Blank"],
            }
            
            if layout_type in layout_names:
                for name in layout_names[layout_type]:
                    for layout in prs.slide_layouts:
                        if name.lower() in layout.name.lower():
                            logger.debug(f"Found layout by name: {layout.name}")
                            return layout
            
            # Fallback to title and content (index 1)
            if len(prs.slide_layouts) > 1:
                logger.warning(f"Layout {layout_type} not found, using fallback")
                return prs.slide_layouts[1]
            
            # Last resort: use first available layout
            if len(prs.slide_layouts) > 0:
                logger.warning(f"Using first available layout as fallback")
                return prs.slide_layouts[0]
            
            return None
            
        except Exception as e:
            logger.error(f"Failed to get layout {layout_type}: {e}")
            return None

    def get_available_layouts(self, prs: Presentation) -> List[Dict[str, Any]]:
        """Get list of available layouts in presentation."""
        layouts = []
        
        for i, layout in enumerate(prs.slide_layouts):
            layout_info = {
                "index": i,
                "name": layout.name,
                "placeholders": []
            }
            
            # Get placeholder information
            for placeholder in layout.placeholders:
                placeholder_info = {
                    "index": placeholder.placeholder_format.idx,
                    "type": placeholder.placeholder_format.type,
                    "name": getattr(placeholder, 'name', 'Unknown')
                }
                layout_info["placeholders"].append(placeholder_info)
            
            layouts.append(layout_info)
        
        return layouts