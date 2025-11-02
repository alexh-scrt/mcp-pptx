"""Theme application to PowerPoint presentations."""

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ..models.theme_spec import ScrapedTheme

logger = logging.getLogger(__name__)


class ThemeApplicator:
    """Applies scraped themes to PowerPoint presentations."""

    def __init__(self):
        self.section_count = 0  # Track section slides for color rotation

    def apply_theme(self, prs: Presentation, theme: ScrapedTheme) -> None:
        """Apply scraped theme to presentation."""
        try:
            self._apply_color_scheme(prs, theme)
            self._apply_fonts(prs, theme)
            # Logo application would be done per slide
            logger.info("Applied theme to presentation")
        except Exception as e:
            logger.error(f"Failed to apply theme: {e}")

    def _apply_color_scheme(self, prs: Presentation, theme: ScrapedTheme) -> None:
        """Apply color scheme to presentation."""
        try:
            # Color scheme will be applied to individual slides
            # via apply_slide_background() method
            logger.debug("Color scheme ready for slide application")

        except Exception as e:
            logger.error(f"Failed to apply color scheme: {e}")

    def apply_slide_background(self, slide, theme: ScrapedTheme, layout_type: str = "TITLE_CONTENT") -> None:
        """Apply background color to a specific slide based on layout type."""
        try:
            # Determine background color based on layout type
            if layout_type == "TITLE":
                # Title slides: RED background
                background_color = self.hex_to_rgb(theme.colors.primary)
            elif layout_type == "SECTION":
                # Section slides: Rotating colors (Cyan, Red, Light Peach)
                section_colors = [
                    theme.colors.accent,      # Cyan #1CCBD0
                    theme.colors.primary,     # Red #E3342F
                    theme.colors.secondary    # Light Peach #FFE9D3
                ]
                color_index = self.section_count % len(section_colors)
                background_color = self.hex_to_rgb(section_colors[color_index])
                self.section_count += 1
                logger.debug(f"Section slide #{self.section_count} using color {section_colors[color_index]}")
            elif layout_type == "CODE":
                # Code slides: WHITE background (code box will have light gray)
                background_color = self.hex_to_rgb(theme.colors.background)
            else:
                # Content slides: WHITE background
                background_color = self.hex_to_rgb(theme.colors.background)

            # Access the slide background
            background = slide.background

            # Set to solid fill
            fill = background.fill
            fill.solid()

            # Apply the background color
            fill.fore_color.rgb = background_color

            logger.debug(f"Applied {layout_type} background color to slide")

        except Exception as e:
            logger.error(f"Failed to apply slide background: {e}")

    def add_title_bar_to_content_slide(self, slide, theme: ScrapedTheme) -> None:
        """Add a colored title bar at the top of content slides (like sample.py)."""
        try:
            # Create title bar shape (rectangle at top of slide)
            # Position: x=0, y=0.4", width=10", height=0.8"
            left = Inches(0)
            top = Inches(0.4)
            width = Inches(10)
            height = Inches(0.8)

            # Add rectangle shape (MSO_SHAPE.RECTANGLE = 1)
            title_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )

            # Fill with primary color (RED)
            title_bar_fill = title_bar.fill
            title_bar_fill.solid()
            title_bar_fill.fore_color.rgb = self.hex_to_rgb(theme.colors.primary)

            # Remove border
            title_bar.line.fill.background()

            logger.debug("Added red title bar to content slide")

        except Exception as e:
            logger.error(f"Failed to add title bar: {e}")

    def _apply_fonts(self, prs: Presentation, theme: ScrapedTheme) -> None:
        """Apply font scheme to presentation."""
        try:
            # This is a simplified implementation
            # In a full implementation, you would modify the theme fonts
            # in the presentation's theme part
            
            logger.debug(f"Fonts will be applied: heading={theme.fonts.heading}, body={theme.fonts.body}")
            
        except Exception as e:
            logger.error(f"Failed to apply fonts: {e}")

    def hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor."""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) != 6:
            hex_color = '000000'  # Default to black
        
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        except ValueError:
            return RGBColor(0, 0, 0)  # Default to black

    def apply_logo_to_slide(self, slide, theme: ScrapedTheme, position: str = "top-right") -> bool:
        """Apply logo to a specific slide."""
        if not theme.logo or not theme.logo.cached_path:
            return False
        
        try:
            logo_path = Path(theme.logo.cached_path)
            if not logo_path.exists():
                logger.warning(f"Logo file not found: {logo_path}")
                return False
            
            # Add logo image to slide
            if position == "top-right":
                left = Inches(8.5)  # Right side
                top = Inches(0.5)   # Top
                width = Inches(1.5)  # Small size
            elif position == "top-left":
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(1.5)
            else:  # center
                left = Inches(4.25)
                top = Inches(3.5)
                width = Inches(2.0)
            
            slide.shapes.add_picture(
                str(logo_path),
                left,
                top,
                width=width
            )
            
            logger.debug(f"Added logo to slide at {position}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to add logo to slide: {e}")
            return False