"""Main PowerPoint presentation renderer."""

import json
import logging
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from ..models.deck_spec import DeckSpec, LayoutType
from ..models.theme_spec import ScrapedTheme, ColorPalette, FontPalette
from .layouts import LayoutManager
from .theme_applicator import ThemeApplicator
from .content_fillers import ContentFiller

logger = logging.getLogger(__name__)


class PresentationRenderer:
    """Renders PowerPoint presentations from DeckSpec."""

    def __init__(self) -> None:
        self.layout_manager = LayoutManager()
        self.theme_applicator = ThemeApplicator()
        self.content_filler = ContentFiller()
        self._default_theme = self._load_default_theme()

    def _load_default_theme(self) -> Optional[ScrapedTheme]:
        """Load the default Secret AI theme."""
        try:
            # Try multiple paths to find the theme file
            possible_paths = [
                Path(__file__).parent.parent.parent.parent / "themes" / "default_theme.json",  # Project root
                Path("themes") / "default_theme.json",  # Relative to cwd
                Path(__file__).parent.parent.parent / "themes" / "default_theme.json"  # src/themes
            ]

            default_theme_path = None
            for path in possible_paths:
                if path.exists():
                    default_theme_path = path
                    break

            if default_theme_path:
                with open(default_theme_path, 'r') as f:
                    theme_data = json.load(f)

                # Create ScrapedTheme from JSON data
                theme = ScrapedTheme(
                    colors=ColorPalette(**theme_data['colors']),
                    fonts=FontPalette(**theme_data['fonts']),
                    logo=None,
                    source_url="default",
                    warnings=[]
                )
                logger.info(f"Loaded default Secret AI theme from {default_theme_path}")
                return theme
        except Exception as e:
            logger.warning(f"Could not load default theme: {e}")

        return None

    async def list_templates(self) -> List[Dict[str, Any]]:
        """List available PowerPoint templates."""
        templates_dir = Path("themes")
        templates = []

        # Add default Secret AI theme
        if self._default_theme:
            default_theme_template = {
                "name": "secret_ai_default",
                "type": "theme",
                "path": "themes/default_theme.json",
                "layouts": [layout.value for layout in LayoutType],
                "description": "Default Secret AI theme (Red, Cyan, White)",
                "colors": {
                    "primary": self._default_theme.colors.primary,
                    "secondary": self._default_theme.colors.secondary,
                    "accent": self._default_theme.colors.accent,
                    "background": self._default_theme.colors.background,
                    "text": self._default_theme.colors.text
                },
                "fonts": {
                    "heading": self._default_theme.fonts.heading,
                    "body": self._default_theme.fonts.body
                }
            }
            templates.append(default_theme_template)

        # Scan for .potx template files
        if templates_dir.exists():
            for template_file in templates_dir.glob("*.potx"):
                template_info = {
                    "name": template_file.stem,
                    "type": "potx",
                    "path": str(template_file),
                    "layouts": [layout.value for layout in LayoutType],
                    "description": f"PowerPoint template: {template_file.stem}"
                }
                templates.append(template_info)

        return templates

    async def generate_presentation(self, deck_spec: DeckSpec) -> Dict[str, Any]:
        """Generate PowerPoint presentation from deck specification."""
        warnings = []
        assets_downloaded = []
        
        try:
            # Create presentation
            if deck_spec.theme.template and Path(deck_spec.theme.template).exists():
                prs = Presentation(deck_spec.theme.template)
                logger.info(f"Using template: {deck_spec.theme.template}")
            else:
                prs = Presentation()
                logger.info("Using default presentation template")
                if deck_spec.theme.template:
                    warnings.append(f"Template {deck_spec.theme.template} not found, using default")
            
            # Apply theme if scraped theme is available, otherwise use default
            theme_to_apply = deck_spec.theme.scraped
            if not theme_to_apply and self._default_theme:
                theme_to_apply = self._default_theme
                logger.info("No theme provided, using default Secret AI theme")

            if theme_to_apply:
                self.theme_applicator.apply_theme(prs, theme_to_apply)
                # Reset section counter for new presentation
                self.theme_applicator.section_count = 0
                logger.info("Applied theme to presentation")
            
            # Generate slides
            slides_generated = 0
            for slide_spec in deck_spec.slides:
                try:
                    # Get appropriate layout
                    layout = self.layout_manager.get_layout(prs, slide_spec.layout)
                    if not layout:
                        layout = self.layout_manager.get_layout(prs, LayoutType.TITLE_CONTENT)
                        warnings.append(f"Layout {slide_spec.layout} not found, using TITLE_CONTENT")
                    
                    # Add slide
                    slide = prs.slides.add_slide(layout)

                    # Apply background color if theme is available
                    if theme_to_apply:
                        # Pass layout type to apply correct background
                        self.theme_applicator.apply_slide_background(
                            slide,
                            theme_to_apply,
                            slide_spec.layout.value
                        )

                        # Add title bar for content slides (TITLE_CONTENT, TWO_COL, etc.)
                        if slide_spec.layout.value in ["TITLE_CONTENT", "TWO_COL", "TABLE", "CHART", "IMAGE_FOCUS"]:
                            self.theme_applicator.add_title_bar_to_content_slide(slide, theme_to_apply)

                    # Fill content
                    slide_warnings = await self.content_filler.fill_slide(
                        slide, slide_spec, theme_to_apply
                    )
                    warnings.extend(slide_warnings)

                    slides_generated += 1

                    # Add footer and slide number
                    self._add_footer_and_slide_number(
                        slide,
                        slides_generated,
                        deck_spec.footer,
                        theme_to_apply
                    )

                    logger.debug(f"Generated slide {slides_generated}: {slide_spec.title}")
                    
                except Exception as e:
                    logger.error(f"Failed to generate slide {slides_generated + 1}: {e}")
                    warnings.append(f"Failed to generate slide {slides_generated + 1}: {str(e)}")

            # Generate output filename
            output_file = self._generate_output_path(deck_spec)
            
            # Ensure output directory exists
            output_file.parent.mkdir(parents=True, exist_ok=True)
            
            # Save presentation
            prs.save(str(output_file))
            logger.info(f"Saved presentation to: {output_file}")
            
            return {
                "ok": True,
                "output": str(output_file),
                "slides_generated": slides_generated,
                "warnings": warnings,
                "assets_downloaded": assets_downloaded
            }
            
        except Exception as e:
            logger.exception("Failed to generate presentation")
            return {
                "ok": False,
                "error": str(e),
                "output": None,
                "slides_generated": 0,
                "warnings": warnings,
                "assets_downloaded": assets_downloaded
            }

    def _generate_output_path(self, deck_spec: DeckSpec) -> Path:
        """Generate output file path."""
        output_dir = Path(deck_spec.output.directory)
        
        if deck_spec.output.filename:
            filename = deck_spec.output.filename
        else:
            # Generate filename from title and timestamp
            safe_title = "".join(c for c in deck_spec.title if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_title = safe_title.replace(' ', '_').lower()
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{safe_title}_{timestamp}.{deck_spec.output.format}"
        
        return output_dir / filename

    def _add_footer_and_slide_number(self, slide, slide_num: int, footer_spec, theme) -> None:
        """Add footer text and slide number to a single slide."""
        try:
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            # Determine footer text
            footer_text = "Generated by Scrt Labs AI"  # Default
            show_slide_number = True  # Default

            if footer_spec:
                if footer_spec.text:
                    footer_text = footer_spec.text
                show_slide_number = footer_spec.show_slide_numbers

            # Determine text color (use theme text color or white as fallback)
            if theme:
                text_color = self.theme_applicator.hex_to_rgb(theme.colors.text)
            else:
                text_color = RGBColor(255, 255, 255)  # White

            # Add footer text box (left side)
            left = Inches(0.5)
            top = Inches(7.0)
            width = Inches(8)
            height = Inches(0.4)

            footer_box = slide.shapes.add_textbox(left, top, width, height)
            footer_frame = footer_box.text_frame
            footer_frame.text = footer_text

            footer_para = footer_frame.paragraphs[0]
            footer_para.font.name = "Tahoma"
            footer_para.font.size = Pt(10)
            footer_para.font.color.rgb = text_color

            # Add slide number (right side) if enabled
            if show_slide_number:
                left = Inches(9)
                top = Inches(7.0)
                width = Inches(0.5)
                height = Inches(0.4)

                number_box = slide.shapes.add_textbox(left, top, width, height)
                number_frame = number_box.text_frame
                number_frame.text = str(slide_num)

                number_para = number_frame.paragraphs[0]
                number_para.alignment = PP_ALIGN.RIGHT
                number_para.font.name = "Tahoma"
                number_para.font.size = Pt(10)
                number_para.font.color.rgb = text_color

            logger.debug(f"Added footer and slide number to slide {slide_num}")

        except Exception as e:
            logger.error(f"Failed to add footer/slide number: {e}")

    def _apply_footer(self, prs: Presentation, footer_spec) -> None:
        """Apply footer to all slides (legacy method, now handled per-slide)."""
        # Footer is now applied per-slide in _add_footer_and_slide_number
        pass