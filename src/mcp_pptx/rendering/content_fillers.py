"""Content filling for PowerPoint slides."""

import logging
from pathlib import Path
from typing import List, Optional

from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..models.deck_spec import SlideSpec, ContentType, ContentPosition
from ..models.theme_spec import ScrapedTheme
from .theme_applicator import ThemeApplicator

logger = logging.getLogger(__name__)


class ContentFiller:
    """Fills slide content based on specifications."""

    def __init__(self) -> None:
        self.theme_applicator = ThemeApplicator()

    async def fill_slide(
        self,
        slide: Slide,
        slide_spec: SlideSpec,
        theme: Optional[ScrapedTheme] = None
    ) -> List[str]:
        """Fill slide with content based on specification."""
        warnings = []

        try:
            # Extract title and subtitle from content array if they have position markers
            title_content = None
            subtitle_content = None
            body_content = []

            for content in slide_spec.content:
                if content.position == ContentPosition.TITLE:
                    title_content = content.text
                elif content.position == ContentPosition.SUBTITLE:
                    subtitle_content = content.text
                else:
                    # Body content or no position specified
                    body_content.append(content)

            # Fill title (from SlideSpec.title or position-based content)
            title_text = slide_spec.title or title_content
            if title_text:
                success = self._fill_title(slide, title_text, theme)
                if not success:
                    warnings.append(f"Could not set title: {title_text}")

            # Fill subtitle (from SlideSpec.subtitle or position-based content)
            subtitle_text = slide_spec.subtitle or subtitle_content
            if subtitle_text:
                success = self._fill_subtitle(slide, subtitle_text, theme)
                if not success:
                    warnings.append(f"Could not set subtitle: {subtitle_text}")

            # Smart content grouping: if we have multiple text-only items, combine them as bullets
            if body_content:
                body_content = self._group_text_items_as_bullets(body_content)

            # Fill body content
            for content in body_content:
                content_warnings = await self._fill_content(slide, content, theme)
                warnings.extend(content_warnings)

            # Add speaker notes
            if slide_spec.speaker_notes:
                self._add_speaker_notes(slide, slide_spec.speaker_notes)

            # Add logo if available
            if theme and theme.logo:
                success = self.theme_applicator.apply_logo_to_slide(slide, theme, "top-right")
                if not success:
                    warnings.append("Could not add logo to slide")

        except Exception as e:
            logger.error(f"Failed to fill slide: {e}")
            warnings.append(f"Error filling slide: {str(e)}")

        return warnings

    def _split_bullet_with_colon(self, bullet: str) -> tuple:
        """
        Split bullet text based on colon rule.
        If bullet has 1-4 words followed by ':', return (bold_part, plain_part).
        Otherwise return (None, None).
        """
        if ':' not in bullet:
            return (None, None)

        # Find the first colon
        colon_index = bullet.index(':')
        prefix = bullet[:colon_index].strip()

        # Count words in prefix (split by whitespace)
        words = prefix.split()

        # Check if it's 1-4 words
        if 1 <= len(words) <= 4:
            # Bold part includes the colon and a space after it
            bold_part = bullet[:colon_index + 1]  # Include the colon
            plain_part = bullet[colon_index + 1:]  # Everything after colon

            # Add space after colon in bold part if there's content after
            if plain_part.strip():
                bold_part += ' '
                plain_part = plain_part.lstrip()

            return (bold_part, plain_part)

        return (None, None)

    def _split_bullet_with_dash(self, bullet: str) -> tuple:
        """
        Split bullet text based on dash rule.
        If bullet has 1-4 words followed by ' - ', return (bold_part, plain_part).
        Otherwise return (None, None).
        """
        if ' - ' not in bullet:
            return (None, None)

        # Find the first ' - '
        dash_index = bullet.index(' - ')
        prefix = bullet[:dash_index].strip()

        # Count words in prefix (split by whitespace)
        words = prefix.split()

        # Check if it's 1-4 words
        if 1 <= len(words) <= 4:
            # Bold part includes the text before ' - '
            bold_part = prefix + ' -'  # Include the dash with one space before
            plain_part = bullet[dash_index + 3:]  # Everything after ' - ' (3 chars: space-dash-space)

            # Add space after dash in bold part if there's content after
            if plain_part.strip():
                bold_part += ' '
                plain_part = plain_part.lstrip()

            return (bold_part, plain_part)

        return (None, None)

    def _split_bullet_for_bold(self, bullet: str) -> tuple:
        """
        Split bullet text into bold and plain parts.
        Tries two rules exclusively:
        1. Colon rule: 1-4 words followed by ':'
        2. Dash rule: 1-4 words followed by ' - '

        Returns (bold_part, plain_part) if a rule matches, otherwise (None, None).
        """
        # Try colon rule first
        bold_part, plain_part = self._split_bullet_with_colon(bullet)
        if bold_part is not None:
            return (bold_part, plain_part)

        # If colon rule doesn't match, try dash rule
        bold_part, plain_part = self._split_bullet_with_dash(bullet)
        if bold_part is not None:
            return (bold_part, plain_part)

        # No rule matched
        return (None, None)

    def _group_text_items_as_bullets(self, content_items: List) -> List:
        """Group consecutive text-only items into a single bullets item."""
        if not content_items:
            return content_items

        # Check if we have multiple text items (common case: plain strings converted to text)
        text_items = []
        other_items = []

        for item in content_items:
            # Check if this is a simple text item (not bullets, not two-column, etc.)
            is_simple_text = (
                item.type == ContentType.TEXT and
                item.text and
                not item.left and
                not item.right
            )
            if is_simple_text:
                text_items.append(item.text)
            else:
                other_items.append(item)

        # If we have multiple text items and no other complex items, convert to bullets
        if len(text_items) > 1 and not other_items:
            # Create a single bullets content item
            from ..models.deck_spec import SlideContent
            bullets_item = SlideContent(
                type=ContentType.BULLETS,
                bullets=text_items
            )
            return [bullets_item]

        # If we have some text items along with other content, still group the text items
        if len(text_items) > 0:
            grouped_items = []
            if text_items:
                from ..models.deck_spec import SlideContent
                # If multiple text items, make them bullets; if single, keep as text
                if len(text_items) > 1:
                    bullets_item = SlideContent(
                        type=ContentType.BULLETS,
                        bullets=text_items
                    )
                    grouped_items.append(bullets_item)
                else:
                    # Single text item, keep as-is
                    grouped_items.append(content_items[0])
            grouped_items.extend(other_items)
            return grouped_items

        # No text items to group, return as-is
        return content_items

    def _fill_title(self, slide: Slide, title: str, theme: Optional[ScrapedTheme] = None) -> bool:
        """Fill slide title."""
        try:
            # Check if this slide has a title bar (content slides with red bar on top)
            # We detect this by looking for a shape at position 0, 0.4" (the title bar)
            has_title_bar = False
            for shape in slide.shapes:
                if (hasattr(shape, 'top') and hasattr(shape, 'left') and
                    shape.top == Inches(0.4) and shape.left == Inches(0)):
                    has_title_bar = True
                    break

            if has_title_bar:
                # Content slide with title bar - create custom positioned title
                # Position: x=0.5", y=0.45", width=9", height=0.7" (on top of red bar)
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.45), Inches(9), Inches(0.7)
                )
                title_frame = title_box.text_frame
                title_frame.text = title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(32)
                title_para.font.bold = True
                title_para.font.name = theme.fonts.heading if theme else "Calibri"
                # White text on red bar
                title_para.font.color.rgb = self.theme_applicator.hex_to_rgb("#FFFFFF")
                return True

            # Find title placeholder (for TITLE and SECTION slides)
            title_placeholder = None
            for shape in slide.placeholders:
                if hasattr(shape, 'text') and 'title' in shape.name.lower():
                    title_placeholder = shape
                    break

            if not title_placeholder:
                # Try placeholder index 0 (usually title)
                if len(slide.placeholders) > 0:
                    title_placeholder = slide.placeholders[0]

            if title_placeholder and hasattr(title_placeholder, 'text'):
                title_placeholder.text = title

                # Apply theme colors if available
                if theme:
                    try:
                        text_frame = title_placeholder.text_frame
                        paragraph = text_frame.paragraphs[0]
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()

                        # Set font
                        run.font.name = theme.fonts.heading
                        run.font.size = Pt(40)  # Large title font (40pt minimum)
                        run.font.bold = True  # Make title bold

                        # Set color to white (for colored backgrounds)
                        white_color = self.theme_applicator.hex_to_rgb("#FFFFFF")
                        run.font.color.rgb = white_color

                    except Exception as e:
                        logger.debug(f"Could not apply theme to title: {e}")

                return True

        except Exception as e:
            logger.error(f"Failed to fill title: {e}")

        return False

    def _fill_subtitle(self, slide: Slide, subtitle: str, theme: Optional[ScrapedTheme] = None) -> bool:
        """Fill slide subtitle."""
        try:
            # Find subtitle placeholder
            subtitle_placeholder = None
            for shape in slide.placeholders:
                if hasattr(shape, 'text') and 'subtitle' in shape.name.lower():
                    subtitle_placeholder = shape
                    break
            
            if not subtitle_placeholder and len(slide.placeholders) > 1:
                # Try placeholder index 1 (usually subtitle)
                subtitle_placeholder = slide.placeholders[1]
            
            if subtitle_placeholder and hasattr(subtitle_placeholder, 'text'):
                subtitle_placeholder.text = subtitle
                
                # Apply theme if available
                if theme:
                    try:
                        text_frame = subtitle_placeholder.text_frame
                        paragraph = text_frame.paragraphs[0]
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()

                        run.font.name = theme.fonts.body
                        run.font.size = Pt(24)  # Subtitle font (24pt for readability)

                        # Use secondary color (light peach) for subtitles - looks good on red backgrounds
                        subtitle_color = self.theme_applicator.hex_to_rgb(theme.colors.secondary)
                        run.font.color.rgb = subtitle_color

                    except Exception as e:
                        logger.debug(f"Could not apply theme to subtitle: {e}")
                
                return True
                
        except Exception as e:
            logger.error(f"Failed to fill subtitle: {e}")
        
        return False

    async def _fill_content(
        self,
        slide: Slide,
        content,
        theme: Optional[ScrapedTheme] = None
    ) -> List[str]:
        """Fill slide content based on content type."""
        warnings = []

        try:
            # Handle two-column content (TWO_COL layout)
            if hasattr(content, 'left') and hasattr(content, 'right') and (content.left or content.right):
                success = self._fill_two_column(slide, content.left or [], content.right or [], theme)
                if not success:
                    warnings.append("Could not add two-column content")
                return warnings

            if content.type == ContentType.TEXT:
                # Check if we have actual text content
                if content.text:
                    success = self._fill_text(slide, content.text, theme)
                    if not success:
                        warnings.append("Could not add text content")
                # If no text but we have left/right (fallback), try two-column
                elif hasattr(content, 'left') or hasattr(content, 'right'):
                    success = self._fill_two_column(
                        slide,
                        getattr(content, 'left', []) or [],
                        getattr(content, 'right', []) or [],
                        theme
                    )
                    if not success:
                        warnings.append("Could not add two-column content")

            elif content.type == ContentType.BULLETS:
                success = self._fill_bullets(slide, content.bullets or [], theme)
                if not success:
                    warnings.append("Could not add bullet points")
                    
            elif content.type == ContentType.IMAGE:
                if content.image:
                    success = await self._fill_image(slide, content.image, theme)
                    if not success:
                        warnings.append(f"Could not add image: {content.image.url}")
                        
            elif content.type == ContentType.TABLE:
                if content.table:
                    success = self._fill_table(slide, content.table, theme)
                    if not success:
                        warnings.append("Could not add table")
                        
            elif content.type == ContentType.CHART:
                if content.chart:
                    success = self._fill_chart(slide, content.chart, theme)
                    if not success:
                        warnings.append("Could not add chart")

            elif content.type == ContentType.CODE:
                if content.code:
                    success = self._fill_code(slide, content.code, theme)
                    if not success:
                        warnings.append("Could not add code content")

        except Exception as e:
            logger.error(f"Failed to fill content: {e}")
            warnings.append(f"Error filling content: {str(e)}")
        
        return warnings

    def _fill_text(self, slide: Slide, text: str, theme: Optional[ScrapedTheme] = None) -> bool:
        """Fill text content."""
        try:
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if hasattr(shape, 'text') and ('content' in shape.name.lower() or 'body' in shape.name.lower()):
                    content_placeholder = shape
                    break
            
            if content_placeholder and hasattr(content_placeholder, 'text'):
                content_placeholder.text = text
                
                # Apply theme
                if theme:
                    try:
                        text_frame = content_placeholder.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = theme.fonts.body
                                run.font.size = Pt(24)  # Body text (24pt for readability)
                                text_color = self.theme_applicator.hex_to_rgb(theme.colors.text)
                                run.font.color.rgb = text_color
                    except Exception as e:
                        logger.debug(f"Could not apply theme to text: {e}")
                
                return True
                
        except Exception as e:
            logger.error(f"Failed to fill text: {e}")
        
        return False

    def _fill_bullets(self, slide: Slide, bullets: List[str], theme: Optional[ScrapedTheme] = None) -> bool:
        """Fill bullet points."""
        try:
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if hasattr(shape, 'text') and ('content' in shape.name.lower() or 'body' in shape.name.lower()):
                    content_placeholder = shape
                    break

            if content_placeholder and hasattr(content_placeholder, 'text_frame'):
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Clear existing content

                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    # Check if bullet has pattern: up to 4 words followed by colon
                    bold_part, plain_part = self._split_bullet_for_bold(bullet)

                    if bold_part and plain_part:
                        # Add bold part
                        run_bold = p.add_run()
                        run_bold.text = bold_part
                        run_bold.font.bold = True

                        # Add plain part
                        run_plain = p.add_run()
                        run_plain.text = plain_part
                    else:
                        # No special formatting needed
                        p.text = bullet

                    p.level = 0  # Top level bullet

                    # Apply theme
                    if theme:
                        try:
                            for run in p.runs:
                                run.font.name = theme.fonts.body
                                run.font.size = Pt(24)  # Bullet font (24pt minimum)
                                text_color = self.theme_applicator.hex_to_rgb(theme.colors.text)
                                run.font.color.rgb = text_color
                        except Exception as e:
                            logger.debug(f"Could not apply theme to bullets: {e}")

                return True

        except Exception as e:
            logger.error(f"Failed to fill bullets: {e}")

        return False

    def _fill_two_column(
        self,
        slide: Slide,
        left_items: List[str],
        right_items: List[str],
        theme: Optional[ScrapedTheme] = None
    ) -> bool:
        """Fill two-column content layout."""
        try:
            # For two-column layouts, we'll add text boxes manually
            # since standard placeholders don't support two columns well

            # Use standard PowerPoint slide dimensions (10" x 7.5")
            slide_width = Inches(10)
            slide_height = Inches(7.5)

            # Calculate column positions (leave margins)
            margin_left = Inches(0.5)
            margin_right = Inches(0.5)
            margin_top = Inches(1.5)  # Below title
            margin_bottom = Inches(0.5)

            column_width = (slide_width - margin_left - margin_right - Inches(0.5)) / 2
            column_height = slide_height - margin_top - margin_bottom

            left_x = margin_left
            right_x = margin_left + column_width + Inches(0.5)

            # Add left column text box
            if left_items:
                left_box = slide.shapes.add_textbox(
                    left_x, margin_top,
                    column_width, column_height
                )
                left_frame = left_box.text_frame
                left_frame.word_wrap = True

                for i, item in enumerate(left_items):
                    if i == 0:
                        p = left_frame.paragraphs[0]
                    else:
                        p = left_frame.add_paragraph()

                    # Check if item has pattern: up to 4 words followed by colon
                    bold_part, plain_part = self._split_bullet_for_bold(item)

                    if bold_part and plain_part:
                        # Add bullet and bold part
                        run_bullet = p.add_run()
                        run_bullet.text = "• "

                        run_bold = p.add_run()
                        run_bold.text = bold_part
                        run_bold.font.bold = True

                        # Add plain part
                        run_plain = p.add_run()
                        run_plain.text = plain_part
                    else:
                        # No special formatting needed
                        p.text = f"• {item}"

                    p.level = 0

                    if theme:
                        try:
                            for run in p.runs:
                                run.font.name = theme.fonts.body
                                run.font.size = Pt(20)  # Two-column content (slightly smaller than bullets)
                                text_color = self.theme_applicator.hex_to_rgb(theme.colors.text)
                                run.font.color.rgb = text_color
                        except Exception as e:
                            logger.debug(f"Could not apply theme to left column: {e}")

            # Add right column text box
            if right_items:
                right_box = slide.shapes.add_textbox(
                    right_x, margin_top,
                    column_width, column_height
                )
                right_frame = right_box.text_frame
                right_frame.word_wrap = True

                for i, item in enumerate(right_items):
                    if i == 0:
                        p = right_frame.paragraphs[0]
                    else:
                        p = right_frame.add_paragraph()

                    # Check if item has pattern: up to 4 words followed by colon
                    bold_part, plain_part = self._split_bullet_for_bold(item)

                    if bold_part and plain_part:
                        # Add bullet and bold part
                        run_bullet = p.add_run()
                        run_bullet.text = "• "

                        run_bold = p.add_run()
                        run_bold.text = bold_part
                        run_bold.font.bold = True

                        # Add plain part
                        run_plain = p.add_run()
                        run_plain.text = plain_part
                    else:
                        # No special formatting needed
                        p.text = f"• {item}"

                    p.level = 0

                    if theme:
                        try:
                            for run in p.runs:
                                run.font.name = theme.fonts.body
                                run.font.size = Pt(20)  # Two-column content (slightly smaller than bullets)
                                text_color = self.theme_applicator.hex_to_rgb(theme.colors.text)
                                run.font.color.rgb = text_color
                        except Exception as e:
                            logger.debug(f"Could not apply theme to right column: {e}")

            return True

        except Exception as e:
            logger.error(f"Failed to fill two-column content: {e}")
            return False

    async def _fill_image(self, slide: Slide, image_spec, theme: Optional[ScrapedTheme] = None) -> bool:
        """Fill image content."""
        try:
            # For now, add placeholder text indicating where image would go
            # In a full implementation, you would download and insert the actual image
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if hasattr(shape, 'text') and ('content' in shape.name.lower() or 'body' in shape.name.lower()):
                    content_placeholder = shape
                    break
            
            if content_placeholder and hasattr(content_placeholder, 'text'):
                placeholder_text = f"[IMAGE: {image_spec.url}]"
                if image_spec.alt_text:
                    placeholder_text += f"\nAlt text: {image_spec.alt_text}"
                if image_spec.caption:
                    placeholder_text += f"\nCaption: {image_spec.caption}"
                
                content_placeholder.text = placeholder_text
                return True
                
        except Exception as e:
            logger.error(f"Failed to fill image: {e}")
        
        return False

    def _fill_table(self, slide: Slide, table_spec, theme: Optional[ScrapedTheme] = None) -> bool:
        """Fill table content."""
        try:
            # For now, add placeholder text
            # In a full implementation, you would create an actual table
            
            content_placeholder = None
            for shape in slide.placeholders:
                if hasattr(shape, 'text') and ('content' in shape.name.lower() or 'body' in shape.name.lower()):
                    content_placeholder = shape
                    break
            
            if content_placeholder and hasattr(content_placeholder, 'text'):
                table_text = f"[TABLE]\nHeaders: {', '.join(table_spec.headers)}\n"
                table_text += f"Rows: {len(table_spec.rows)} rows of data"
                
                content_placeholder.text = table_text
                return True
                
        except Exception as e:
            logger.error(f"Failed to fill table: {e}")
        
        return False

    def _fill_chart(self, slide: Slide, chart_spec, theme: Optional[ScrapedTheme] = None) -> bool:
        """Fill chart content."""
        try:
            # For now, add placeholder text
            # In a full implementation, you would create an actual chart

            content_placeholder = None
            for shape in slide.placeholders:
                if hasattr(shape, 'text') and ('content' in shape.name.lower() or 'body' in shape.name.lower()):
                    content_placeholder = shape
                    break

            if content_placeholder and hasattr(content_placeholder, 'text'):
                chart_text = f"[CHART: {chart_spec.type}]"
                if chart_spec.title:
                    chart_text += f"\nTitle: {chart_spec.title}"

                content_placeholder.text = chart_text
                return True

        except Exception as e:
            logger.error(f"Failed to fill chart: {e}")

        return False

    def _split_code_into_chunks(self, code: str, max_lines_per_slide: int = 25) -> List[str]:
        """Split code into chunks that fit on a slide.

        Args:
            code: The complete code string
            max_lines_per_slide: Maximum number of lines per slide (default 25)

        Returns:
            List of code chunks, one per slide
        """
        lines = code.split('\n')

        if len(lines) <= max_lines_per_slide:
            return [code]

        chunks = []
        current_chunk = []

        for line in lines:
            current_chunk.append(line)

            if len(current_chunk) >= max_lines_per_slide:
                chunks.append('\n'.join(current_chunk))
                current_chunk = []

        # Add remaining lines if any
        if current_chunk:
            chunks.append('\n'.join(current_chunk))

        return chunks

    def _fill_code(self, slide: Slide, code_input, theme: Optional[ScrapedTheme] = None) -> bool:
        """Fill code content with Courier New font and light gray background."""
        try:
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            # Parse code input (can be CodeSpec or plain string)
            if isinstance(code_input, str):
                code_text = code_input
                code_title = None
                language = None
            else:
                # It's a CodeSpec object
                code_text = code_input.code
                code_title = code_input.title if hasattr(code_input, 'title') else None
                language = code_input.language if hasattr(code_input, 'language') else None

            # Calculate dimensions for code block
            # Leave room for title bar if present
            margin_left = Inches(0.5)
            margin_top = Inches(1.3) if code_title else Inches(0.5)
            margin_right = Inches(0.5)
            margin_bottom = Inches(0.5)

            code_width = Inches(9)  # 10" - 0.5" margins on each side
            code_height = Inches(6.5) if code_title else Inches(7)  # Adjust for title

            # Add title if present
            if code_title:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5),
                    Inches(9), Inches(0.6)
                )
                title_frame = title_box.text_frame
                title_frame.text = code_title
                title_para = title_frame.paragraphs[0]
                title_para.font.name = "Calibri"
                title_para.font.size = Pt(28)
                title_para.font.bold = True
                if theme:
                    title_para.font.color.rgb = self.theme_applicator.hex_to_rgb(theme.colors.text)
                else:
                    title_para.font.color.rgb = RGBColor(0, 0, 0)

            # Add code text box with light gray background
            code_box = slide.shapes.add_textbox(
                margin_left, margin_top,
                code_width, code_height
            )

            # Set light gray background (#F5F5F5)
            code_box.fill.solid()
            code_box.fill.fore_color.rgb = RGBColor(245, 245, 245)

            # Remove border
            code_box.line.fill.background()

            # Add code text
            code_frame = code_box.text_frame
            code_frame.word_wrap = True
            code_frame.text = code_text

            # Format code text
            for paragraph in code_frame.paragraphs:
                paragraph.font.name = "Courier New"
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                paragraph.alignment = PP_ALIGN.LEFT

            logger.debug(f"Added code block to slide (language: {language})")
            return True

        except Exception as e:
            logger.error(f"Failed to fill code: {e}")
            import traceback
            logger.error(traceback.format_exc())

        return False

    def _add_speaker_notes(self, slide: Slide, notes: str) -> None:
        """Add speaker notes to slide."""
        try:
            if hasattr(slide, 'notes_slide') and hasattr(slide.notes_slide, 'notes_text_frame'):
                slide.notes_slide.notes_text_frame.text = notes
                logger.debug("Added speaker notes to slide")
        except Exception as e:
            logger.error(f"Failed to add speaker notes: {e}")