from pptx import Presentation
from pptx.slide import SlideLayout, Slide, SlideMaster, Slides
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Secret AI color scheme
PRIMARY = RGBColor(227, 52, 47)      # #E3342F - Red
SECONDARY = RGBColor(255, 233, 211)  # #FFE9D3 - Light peach
ACCENT = RGBColor(28, 203, 208)      # #1CCBD0 - Cyan
BACKGROUND = RGBColor(255, 255, 255) # #FFFFFF - White
TEXT = RGBColor(17, 24, 39)          # #111827 - Dark gray

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slides: Slides = prs.slides
    slide: Slide = slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = BACKGROUND
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = SECONDARY
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = BACKGROUND
    title_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND
    
    # Title with accent bar
    title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0.4), Inches(10), Inches(0.8))
    title_bar_fill = title_bar.fill
    title_bar_fill.solid()
    title_bar_fill.fore_color.rgb = PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BACKGROUND
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.space_after = Pt(12)

# Slide 1: Title
add_title_slide(prs, "Confidential AI Infrastructure", 
                "Intel TDX + NVIDIA H100 CC Deployment Guide for DevOps")

# Slide 2: Section 1
add_section_slide(prs, "Section 1: Architecture Overview")

# Slide 3: What is Confidential Computing
add_content_slide(prs, "What is Confidential Computing?", [
    "Hardware-enforced protection for data in use",
    "Memory encrypted at CPU level - inaccessible to host/hypervisor/other VMs",
    "Trust Domains (TDs): Isolated VMs with ephemeral per-VM encryption keys",
    "Remote attestation: Cryptographic proof of VM integrity before use",
    "Zero-trust cloud: Run sensitive AI workloads on untrusted infrastructure",
    "Compliance: GDPR, HIPAA, data sovereignty requirements"
])

# Slide 4: Full Stack Architecture
add_content_slide(prs, "Full Stack Architecture", [
    "Layer 6: Guest VM (Ubuntu 24.04 + NVIDIA Driver + CUDA 12.5+)",
    "Layer 5: OVMF TDX-aware UEFI firmware (TD initialization + measurements)",
    "Layer 4: QEMU userspace (TDX guest object + VFIO GPU passthrough)",
    "Layer 3: KVM kernel module (SEAMCALL interface + TD lifecycle)",
    "Layer 2: Linux kernel 6.9+ (TDX host support + crypto + IOMMU)",
    "Layer 1: Hardware (Intel Emerald Rapids TDX + H100 CC + TDX-SEAM)"
])

# Slide 5: Intel TDX Security
add_content_slide(prs, "Intel TDX Security Guarantees", [
    "Memory Encryption: All TD memory encrypted with ephemeral keys",
    "CPU State Protection: Register contents hidden from host/VMM",
    "Secure EPT: Page tables managed by TDX-SEAM, prevents memory tampering",
    "PAMT: Physical Address Metadata Table (~8GB overhead) tracks ownership",
    "No host access: Even root cannot read TD memory or inject code",
    "Attestation: TD Quote proves configuration integrity to remote verifiers"
])

# Slide 6: NVIDIA H100 CC
add_content_slide(prs, "NVIDIA H100 Confidential Computing", [
    "CC Mode: Secure channel between guest driver and GPU firmware",
    "SPDM Protocol: Mutual authentication + encrypted communication (ECDH)",
    "Attestation Reports: Cryptographic proof of GPU firmware integrity",
    "RIM Verification: Reference Integrity Manifests validate measurements",
    "Host isolation: Host cannot inspect GPU workloads or memory",
    "One-time setup: Set CC mode on host, driver auto-establishes secure channel"
])

# Slide 7: Section 2
add_section_slide(prs, "Section 2: Hardware & Prerequisites")

# Slide 8: Hardware Requirements
add_content_slide(prs, "Hardware Requirements", [
    "CPU: Intel Emerald Rapids (4th/5th Gen Xeon Scalable) with TDX support",
    "GPU: NVIDIA H100 PCIe/SXM5, firmware 96.00.5E.00.00+",
    "Memory: 128GB+ DDR5 (256GB+ recommended, +8GB for PAMT overhead)",
    "Storage: 500GB+ NVMe SSD for host, 100GB+ per guest VM",
    "Network: 10GbE+ for production, IOMMU-capable NIC preferred",
    "Verification: cpuid | grep -i tdx, lspci -d 10de:"
])

# Slide 9: BIOS Configuration
add_content_slide(prs, "Critical BIOS Configuration", [
    "Total Memory Encryption (TME): ENABLED",
    "TME Multi-Tenant (TME-MT): ENABLED",
    "Intel TDX: ENABLED",
    "TDX SEAM Loader: ENABLED",
    "TDX Key Split: Set to non-zero value (e.g., 64)",
    "Software Guard Extensions (SGX): ENABLED",
    "Limit CPU PA to 46 Bits: DISABLED (critical!)",
    "Verify: dmesg | grep 'virt/tdx: BIOS enabled'"
])

# Slide 10: Software Stack
add_content_slide(prs, "Software Stack Requirements", [
    "Host OS: Ubuntu 24.04 LTS (mandatory - 22.04 NOT supported)",
    "Kernel: Custom-built 6.9+ with TDX patches from Intel",
    "QEMU: TDX-patched version with guest object support",
    "OVMF: TDX-aware UEFI firmware (EDK2 stable202405+)",
    "Guest OS: Ubuntu 24.04/22.04 LTS with TDX guest driver",
    "NVIDIA Driver: r550+ in guest, nvTrust tools on host",
    "CUDA: 12.5+ in guest for AI workloads"
])

# Slide 11: Section 3
add_section_slide(prs, "Section 3: Kernel Build & Configuration")

# Slide 12: Kernel Patching
add_content_slide(prs, "Kernel Patching Process", [
    "Clone Intel TDX kernel tree: github.com/intel/tdx-linux",
    "Checkout device-passthrough branch: commit 1323f7b1ddf",
    "Install dependencies: build-essential, libncurses-dev, bison, flex, libssl-dev",
    "Configure kernel: make menuconfig - enable TDX host/guest, VFIO, IOMMU",
    "Enable crypto: AES-NI, SHA-256, RSA, ECDH for attestation",
    "Build: make -j$(nproc) (30-60 minutes)",
    "Install: make modules_install && make install && update-grub"
])

# Slide 13: Kernel Config
add_content_slide(prs, "Essential Kernel Config Options", [
    "CONFIG_INTEL_TDX_HOST=y - TDX host virtualization support",
    "CONFIG_KVM_INTEL=y - KVM for Intel processors",
    "CONFIG_VFIO_PCI=m - VFIO PCI device assignment",
    "CONFIG_VFIO_IOMMU_TYPE1=m - IOMMU support for VFIO",
    "CONFIG_IOMMUFD=y - New IOMMU subsystem for TDX",
    "CONFIG_CRYPTO_AES_NI_INTEL=y - Hardware crypto acceleration",
    "CONFIG_X86_SGX=y - Software Guard Extensions support"
])

# Slide 14: Boot Parameters
add_content_slide(prs, "Boot Parameters & Verification", [
    "GRUB config: /etc/default/grub",
    "Add: GRUB_CMDLINE_LINUX=\"nohibernate kvm_intel.tdx=on\"",
    "Update GRUB: sudo update-grub && sudo reboot",
    "Verify kernel: uname -r (should show custom version)",
    "Check TDX: dmesg | grep -i tdx (look for 'BIOS enabled' message)",
    "Verify KVM: ls /dev/kvm (should exist with correct permissions)",
    "Test SEAMCALL: TDX module should be loaded and functional"
])

# Slide 15: Section 4
add_section_slide(prs, "Section 4: QEMU & Firmware Build")

# Slide 16: QEMU Build
add_content_slide(prs, "Building TDX-Patched QEMU", [
    "Clone QEMU: git clone https://gitlab.com/qemu-project/qemu",
    "Checkout: commit ff6d8490e33 (TDX support branch)",
    "Configure: ./configure --enable-kvm --target-list=x86_64-softmmu",
    "Build: make -j$(nproc) (15-30 minutes)",
    "Install: sudo make install (installs to /usr/local/bin)",
    "Verify: qemu-system-x86_64 --version (check for TDX support)",
    "Key features: TDX guest object, VFIO enhancements, IOMMUFD support"
])

# Slide 17: OVMF Build
add_content_slide(prs, "Building TDX-Aware OVMF", [
    "Clone EDK2: git clone -b edk2-stable202405 github.com/tianocore/edk2",
    "Initialize submodules: git submodule update --init",
    "Set environment: source edksetup.sh",
    "Build BaseTools: make -C BaseTools",
    "Build OVMF: build -a X64 -t GCC5 -p OvmfPkg/OvmfPkgX64.dsc",
    "Output: Build/OvmfX64/DEBUG_GCC5/FV/OVMF.fd",
    "Purpose: TDX measurement, secure boot, TD initialization"
])

# Slide 18: Section 5
add_section_slide(prs, "Section 5: GPU Configuration")

# Slide 19: GPU CC Mode
add_content_slide(prs, "Enabling GPU Confidential Computing", [
    "Clone nvTrust: git clone github.com/NVIDIA/nvtrust.git",
    "Navigate: cd nvtrust/host_tools/python",
    "Query current state: nvidia_gpu_tools.py --gpu-name=H100 --query-cc-mode",
    "Enable CC mode: --set-cc-mode=on --reset-after-cc-mode-switch",
    "GPU will automatically reset (takes ~30 seconds)",
    "Verify: --query-cc-mode (should return 'CC mode is on')",
    "Note: CC mode persists across host reboots"
])

# Slide 20: GPU Passthrough
add_content_slide(prs, "GPU Passthrough Configuration", [
    "Find GPU PCI address: lspci -d 10de: -nn (e.g., 0000:17:00.0)",
    "Get device ID: lspci -d 10de: -nn (e.g., [10de:2321] for H100)",
    "Load VFIO modules: modprobe vfio vfio_pci",
    "Bind to VFIO: echo '10de 2321' > /sys/bus/pci/drivers/vfio-pci/new_id",
    "Verify binding: lspci -k -s 0000:17:00.0 (should show vfio-pci driver)",
    "Check IOMMU group: ls /sys/kernel/iommu_groups/",
    "Unbind from nvidia driver if necessary before binding to vfio-pci"
])

# Slide 21: Section 6
add_section_slide(prs, "Section 6: VM Deployment")

# Slide 22: Guest Prep
add_content_slide(prs, "Guest VM Preparation", [
    "Download ISO: ubuntu-24.04-live-server-amd64.iso",
    "Create disk: qemu-img create -f qcow2 guest.qcow2 500G",
    "Install OS: Boot with KVM (no TDX) using ISO for initial setup",
    "Configure GRUB in guest: console=ttyS0 clearcpuid=mtrr,avx,avx2",
    "Enable LKCA: For NVIDIA driver compatibility in TDX environment",
    "Update guest kernel: 6.8+ with TDX guest driver support",
    "Prepare for TDX launch: Shutdown after initial configuration"
])

# Slide 23: VM Launch
add_content_slide(prs, "Launching Confidential VM", [
    "Key QEMU parameters:",
    "-object tdx-guest,id=tdx (creates TDX guest object)",
    "-machine q35,confidential-guest-support=tdx,kernel-irqchip=split",
    "-object memory-backend-ram,id=mem0,size=64G",
    "-bios /path/to/OVMF.fd (TDX-aware firmware)",
    "-object iommufd,id=iommufd0 (for GPU isolation)",
    "-device vfio-pci,host=17:00.0,iommufd=iommufd0 (GPU passthrough)",
    "-nographic (serial console access)"
])

# Slide 24: Section 7
add_section_slide(prs, "Section 7: Guest Configuration")

# Slide 25: NVIDIA Driver
add_content_slide(prs, "NVIDIA Driver Installation", [
    "Add CUDA repository: cuda-keyring_1.1-1_all.deb",
    "Install driver: apt install nvidia-driver-550-server-open",
    "Configure persistence: Edit nvidia-persistenced.service",
    "Add flags: --uvm-persistence-mode --verbose",
    "Reload systemd: systemctl daemon-reload",
    "Reboot guest: GPU driver will auto-establish SPDM secure channel",
    "Verify: nvidia-smi (should show H100 detected)"
])

# Slide 26: CUDA & AI
add_content_slide(prs, "CUDA Toolkit & AI Frameworks", [
    "Install CUDA: apt install cuda-toolkit-12-5",
    "Set environment: export PATH=/usr/local/cuda-12.5/bin:$PATH",
    "Verify CUDA: nvcc --version",
    "Install PyTorch: pip3 install torch --index-url pytorch.org/whl/cu125",
    "Install TensorFlow: pip3 install tensorflow[and-cuda]",
    "Test GPU access: python3 -c 'import torch; print(torch.cuda.is_available())'",
    "All computation now hardware-protected in confidential VM"
])

# Slide 27: Section 8
add_section_slide(prs, "Section 8: Attestation & Security")

# Slide 28: Attestation
add_content_slide(prs, "Remote Attestation Process", [
    "Install SDK: pip3 install nv-attestation-sdk",
    "Create client: attestation.Attestation('confidential-vm')",
    "Add verifier: client.add_verifier(Devices.GPU, Environment.LOCAL)",
    "Run attestation: result = client.attest() (returns True/False)",
    "Validate policy: client.validate_token(policy_json)",
    "Set GPU ready: nvidia-smi conf-compute -srs 1",
    "Attestation proves: VM integrity, GPU firmware, secure channel established"
])

# Slide 29: Security Checklist
add_content_slide(prs, "Security Verification Checklist", [
    "TDX detected: dmesg | grep 'tdx: Guest detected'",
    "Memory encrypted: All guest RAM protected by hardware encryption",
    "GPU in CC mode: nvidia-smi conf-compute -f (shows CC status: ON)",
    "SPDM channel: Driver logs show successful SPDM handshake",
    "Attestation passing: GPU measurements match RIM values",
    "VFIO isolation: GPU bound to vfio-pci, IOMMU active",
    "No debug interfaces: Disable kernel debug for production"
])

# Slide 30: Section 9
add_section_slide(prs, "Section 9: Operations & Monitoring")

# Slide 31: Timeline
add_content_slide(prs, "Deployment Timeline", [
    "Phase 1 - Hardware setup & BIOS config: 2-4 hours",
    "Phase 2 - Host OS preparation: 1-2 hours",
    "Phase 3 - Kernel build & installation: 2-3 hours",
    "Phase 4 - QEMU/OVMF build: 2-3 hours (can parallelize)",
    "Phase 5 - GPU CC setup: 30 minutes",
    "Phase 6-10 - VM creation to AI deployment: 3-4 hours",
    "Total critical path: 12-16 hours for first deployment"
])

# Slide 32: Troubleshooting
add_content_slide(prs, "Common Troubleshooting Scenarios", [
    "TDX not detected: Check BIOS settings, especially 'Limit CPU PA to 46 Bits'",
    "SEAMCALL errors: Verify kernel config has CONFIG_INTEL_TDX_HOST=y",
    "GPU passthrough fails: Check vfio-pci binding, IOMMU groups",
    "VM won't boot: Verify OVMF.fd path, check memory backend configuration",
    "nvidia-smi not working: Check clearcpuid parameters, driver installation",
    "Attestation failing: Verify GPU firmware version, check RIM files",
    "Performance issues: Enable hugepages, check NUMA configuration"
])

# Slide 33: Best Practices
add_content_slide(prs, "Production Best Practices", [
    "Automation: Script kernel build, VM provisioning, attestation checks",
    "Monitoring: Set up dmesg logging, GPU metrics, attestation status",
    "Backups: Image guest VMs regularly, document kernel config",
    "Updates: Track Intel TDX kernel updates, NVIDIA driver releases",
    "Security: Regular attestation, disable debug features, audit access",
    "Documentation: Maintain runbooks, configuration baseline, recovery procedures",
    "Testing: Validate on staging before production, disaster recovery drills"
])

# Slide 34: Key Takeaways
add_content_slide(prs, "Key Takeaways", [
    "Hardware-enforced confidentiality: Host/hypervisor cannot access VM memory",
    "Full stack protection: CPU + GPU both in confidential computing mode",
    "Custom kernel required: Ubuntu 24.04 + Intel TDX patches mandatory",
    "BIOS critical: One wrong setting (PA limit) breaks entire deployment",
    "Attestation essential: Always verify before processing sensitive data",
    "Production ready: 12-16 hours initial setup, mature for sensitive AI workloads",
    "DevOps friendly: Automate build/deploy, integrate with CI/CD pipelines"
])

# Slide 35: Resources
add_content_slide(prs, "Resources & Next Steps", [
    "Documentation: Review qemu-kvm-arch.html, tdx-cc-architecture.html",
    "Kernel guide: Complete step-by-step in Linux Kernel Patching Guide",
    "Intel TDX: github.com/intel/tdx-linux (official kernel patches)",
    "NVIDIA nvTrust: github.com/NVIDIA/nvtrust (GPU attestation tools)",
    "Community: Intel TDX documentation, NVIDIA developer forums",
    "Training: 2-day intensive workshop available",
    "Support: Internal Slack #confidential-ai, create repository issues"
])

# Save
prs.save('./confidential_ai_devops_presentation.pptx')
print("Presentation created successfully!")